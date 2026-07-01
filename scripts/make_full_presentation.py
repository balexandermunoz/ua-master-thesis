"""Generate full dissertation PowerPoint presentation.

Covers Chapter 3 - Methodology, Chapter 4 - Implementation & Results,
and Chapter 5 - Conclusions.  (~38 slides, ~20-minute presentation)
"""

from pathlib import Path
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── paths ──────────────────────────────────────────────────────────────
ROOT     = Path(__file__).resolve().parent.parent
OUT      = ROOT / "Dissertation_Presentation.pptx"
ARCH_IMG = ROOT / "figs" / "architecture.png"
UI_IMG   = ROOT / "figs" / "UIM1_light.png"
DEMO_VIDEO = ROOT / "figs" / "shortDemo.mp4"
LOGO_WHITE = ROOT / "figs" / "UbiwhereLogo.png"
LOGO_BLUE = ROOT / "figs" / "UbiwhereLogoBlue.png"
NOTES_FILE = ROOT / "scripts" / "presenter_notes.txt"

# ── colour palette ─────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG    = RGBColor(0x1B, 0x2A, 0x4A)   # dark navy
ACCENT     = RGBColor(0x00, 0x96, 0xD6)   # bright blue
ACCENT2    = RGBColor(0x00, 0xB4, 0x8A)   # teal-green
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
MED_GRAY   = RGBColor(0x6C, 0x75, 0x7D)
BLACK      = RGBColor(0x20, 0x20, 0x20)
ORANGE     = RGBColor(0xFF, 0x8C, 0x00)
GREEN      = RGBColor(0x28, 0xA7, 0x45)
RED_LIGHT  = RGBColor(0xDC, 0x35, 0x45)
LAYER1_CLR = RGBColor(0x00, 0x96, 0xD6)
LAYER2_CLR = RGBColor(0x00, 0xB4, 0x8A)
LAYER3_CLR = RGBColor(0xFF, 0x8C, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── helper functions ───────────────────────────────────────────────────

def _add_footer(slide, slide_num, total, dark=False):
    """Add a footer bar with slide number, company name, and email."""
    y = 7.08
    fg = RGBColor(0xAA, 0xBB, 0xCC) if dark else MED_GRAY
    logo_path = LOGO_WHITE if dark else LOGO_BLUE
    if logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(0.7), Inches(6.98), height=Inches(0.32))
    else:
        txBox = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(3.0), Inches(0.3))
        p = txBox.text_frame.paragraphs[0]
        p.text = "Ubiwhere"
        p.font.size = Pt(10); p.font.color.rgb = fg; p.font.name = "Calibri"
        p.font.bold = True; p.alignment = PP_ALIGN.LEFT
    txBox = slide.shapes.add_textbox(Inches(4.2), Inches(y), Inches(5.0), Inches(0.3))
    p = txBox.text_frame.paragraphs[0]
    p.text = "balexandermunoz@ua.pt"
    p.font.size = Pt(10); p.font.color.rgb = fg; p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    txBox = slide.shapes.add_textbox(Inches(10.3), Inches(y), Inches(2.3), Inches(0.3))
    p = txBox.text_frame.paragraphs[0]
    p.text = str(slide_num)
    p.font.size = Pt(10); p.font.color.rgb = fg; p.font.name = "Calibri"
    p.alignment = PP_ALIGN.RIGHT


def _load_presenter_notes(file_path):
    """Load presenter notes in [Slide N] sections from a plain text file."""
    notes = {}
    if not file_path.exists():
        return notes

    content = file_path.read_text(encoding="utf-8")
    current_slide = None
    buffer = []

    for raw_line in content.splitlines():
        line = raw_line.rstrip("\n")
        match = re.match(r"^\[Slide\s+(\d+)\]\s*$", line.strip(), re.IGNORECASE)
        if match:
            if current_slide is not None:
                notes_text = "\n".join(buffer).strip()
                if notes_text:
                    notes[current_slide] = notes_text
            current_slide = int(match.group(1))
            buffer = []
            continue

        if current_slide is not None:
            buffer.append(line)

    if current_slide is not None:
        notes_text = "\n".join(buffer).strip()
        if notes_text:
            notes[current_slide] = notes_text

    return notes


def _apply_presenter_notes(presentation, notes_by_slide):
    """Write presenter notes into each slide notes pane."""
    for idx, slide in enumerate(presentation.slides, start=1):
        notes_text = notes_by_slide.get(idx)
        if not notes_text:
            continue

        notes_frame = slide.notes_slide.notes_text_frame
        notes_frame.clear()
        notes_frame.paragraphs[0].text = notes_text


def _is_dark_background(slide):
    """Return True when slide background is dark enough for light footer assets."""
    try:
        fill = slide.background.fill
        rgb = fill.fore_color.rgb
        if rgb is None:
            return False
        r, g, b = rgb[0], rgb[1], rgb[2]
        # Perceived luminance (ITU-R BT.709)
        luminance = 0.2126 * r + 0.7152 * g + 0.0722 * b
        return luminance < 128
    except Exception:
        return False

def _solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(font_size * 0.3)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return tf

def _add_bullet_list(slide, left, top, width, height, items,
                     font_size=16, color=BLACK, font_name="Calibri",
                     bullet_color=ACCENT, line_spacing=1.35, bold_prefix=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)
        run_b = p.add_run()
        run_b.text = "●  "
        run_b.font.size = Pt(font_size - 2)
        run_b.font.color.rgb = bullet_color
        run_b.font.name = font_name
        if bold_prefix and ": " in item:
            prefix, rest = item.split(": ", 1)
            run_prefix = p.add_run()
            run_prefix.text = prefix + ": "
            run_prefix.font.size = Pt(font_size)
            run_prefix.font.bold = True
            run_prefix.font.color.rgb = color
            run_prefix.font.name = font_name
            run_rest = p.add_run()
            run_rest.text = rest
            run_rest.font.size = Pt(font_size)
            run_rest.font.color.rgb = color
            run_rest.font.name = font_name
        else:
            run_t = p.add_run()
            run_t.text = item
            run_t.font.size = Pt(font_size)
            run_t.font.color.rgb = color
            run_t.font.name = font_name
    return tf

def _add_rect(slide, left, top, width, height, fill_color, text="",
              font_size=14, font_color=WHITE, bold=True, corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = "Calibri"
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)
    shape.text_frame.auto_size = None
    shape.text_frame.word_wrap = True
    from pptx.enum.text import MSO_ANCHOR
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    try:
        shape.text_frame.anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    return shape

def _slide_title(slide, title_text, subtitle_text=None, dark=False):
    fg = WHITE if dark else DARK_BG
    _add_textbox(slide, 0.7, 0.35, 11.5, 0.7, title_text,
                 font_size=32, bold=True, color=fg)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.05),
        Inches(2.5), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    if subtitle_text:
        _add_textbox(slide, 0.7, 1.2, 11.5, 0.5, subtitle_text,
                     font_size=18, color=MED_GRAY if not dark else RGBColor(0xAA, 0xBB, 0xCC))

def _add_table_slide(slide, headers, rows, left=0.7, top=1.6, width=11.9,
                     col_widths=None, font_size=13, row_height=0.4):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                         Inches(left), Inches(top),
                                         Inches(width), Inches(row_height * n_rows))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BG
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = "Calibri"
            paragraph.alignment = PP_ALIGN.CENTER
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = BLACK
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER
    return table

# ======================================================================
#  CHAPTER 3 - METHODOLOGY (slides 1-16, unchanged)
# ======================================================================

# ── 1. TITLE SLIDE ────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, DARK_BG)
_add_textbox(sl, 1.0, 1.8, 11.3, 1.2,
             "Master's Dissertation Presentation",
             font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
_add_textbox(sl, 1.0, 3.2, 11.3, 0.7,
             "Development of a Modular Simulation Algorithm for Different Environments",
             font_size=22, color=ACCENT, alignment=PP_ALIGN.CENTER)
line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                           Inches(4.5), Inches(4.1), Inches(4.3), Inches(0.06))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
_add_textbox(sl, 1.0, 4.4, 11.3, 0.5,
             "Ubiwhere  ·  University of Aveiro  · April 2026",
             font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ── 2. OUTLINE ────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Outline")

outline_items = [
    "1. 🎯 Motivation & Objective",
    "2. 📚 State of the Art",
    "3. 🧪 Research Approach",
    "4. 🏗️ System Architecture",
    "5. ⚙️ Design Decisions",
    "6. 🧭 Simulation Scenarios",
    "7. 📊 Evaluation Methodology",
    "8. 💻 Implementation",
    "9. 📈 Results",
    "10. ✅ Conclusions",
]

left_x = 1.0
right_x = 6.7
start_y = 1.8
line_gap = 0.55

for i, label in enumerate(outline_items):
    if i < 5:
        x = left_x
        y = start_y + i * line_gap
    else:
        x = right_x
        y = start_y + (i - 5) * line_gap
    _add_textbox(sl, x, y, 4.8, 0.35, label,
                 font_size=17, color=DARK_BG, bold=True)

_add_textbox(sl, 0.7, 6.0, 11.5, 0.4,
             "A concise roadmap of the dissertation narrative: from motivation to conclusions.",
             font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ── 3. MOTIVATION & OBJECTIVE ─────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Motivation & Objective")

_add_textbox(sl, 0.7, 1.6, 5.5, 0.5,
             "The Problem", font_size=22, bold=True, color=DARK_BG)
_add_bullet_list(sl, 0.7, 2.2, 5.5, 3.0, [
    "Cities depend on interdependent infrastructure: energy, mobility, telecom (EV, PowerGrids, Sync)",
    "Existing simulators operate in isolation - one domain each, not integrated",
    "Urban platforms (e.g. Ubiwhere UBP) monitor data but cannot predict cross-domain effects",
], font_size=15, bullet_color=ACCENT, bold_prefix=False, line_spacing=1.5)

_add_textbox(sl, 7.0, 1.6, 5.5, 0.5,
             "Our Goal", font_size=22, bold=True, color=DARK_BG)
_add_bullet_list(sl, 7.0, 2.2, 5.5, 3.0, [
    "Enable urban platforms to answer cross-domain \"what-if\" questions",
    "Build a flexible, reusable, domain-agnostic simulation engine",
    "Support multi-domain co-simulation (energy + mobility + telecom)",
    "Validate through practical scenarios - standalone and coupled",
], font_size=15, bullet_color=ACCENT2, bold_prefix=False, line_spacing=1.5)

_add_rect(sl, 0.7, 5.3, 11.9, 0.7, DARK_BG,
          "Key question:  What system-level effects are invisible when domains are simulated in isolation?",
          font_size=16, font_color=WHITE, bold=True)

# ── 4. STATE OF THE ART RECAP ─────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "State of the Art", "Domain-specific simulators & co-simulation frameworks")

domains_sota = [
    ("Energy Simulators", LAYER1_CLR, [
        "PyPSA - optimization & sector coupling",
        "GridLAB-D - agent-based smart grid",
        "Pandapower - Python, distribution networks",
        "OpenDSS - industry standard, DER integration",
    ]),
    ("Mobility Simulators", LAYER2_CLR, [
        "SUMO - microscopic, TraCI interface",
        "MATSim - large-scale, activity-based",
        "OpenTrafficSim - multi-level (micro/meso/macro)",
    ]),
    ("Telecom Simulators", LAYER3_CLR, [
        "ns-3 - discrete-event, protocol fidelity",
        "SimuLTE/Simu5G - LTE/5G system-level",
        "OMNeT++ - modular, event-driven framework",
    ]),
]
for i, (title, clr, items) in enumerate(domains_sota):
    x = 0.7 + i * 4.2
    _add_rect(sl, x, 1.7, 3.8, 0.6, clr, title, font_size=16)
    _add_bullet_list(sl, x + 0.15, 2.5, 3.5, 2.5, items,
                     font_size=13, bullet_color=clr, bold_prefix=False, line_spacing=1.45)

_add_textbox(sl, 0.7, 4.6, 11.5, 0.4,
             "Co-Simulation Frameworks", font_size=20, bold=True, color=DARK_BG)

cosim_items = [
    ("Mosaik", MED_GRAY, "Smart-grid focus, Python API"),
    ("MECSYCO", MED_GRAY, "DEVS-based, multi-paradigm"),
    ("HELICS  ✓", RGBColor(0xA1, 0x6C, 0xFF), "Multi-domain, scalable, flexible time mgmt"),
]
for i, (name, clr, desc) in enumerate(cosim_items):
    x = 0.7 + i * 4.2
    _add_rect(sl, x, 5.15, 3.8, 0.45, clr, name, font_size=15)
    _add_textbox(sl, x, 5.7, 3.8, 0.4, desc,
                 font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

_add_textbox(sl, 0.7, 6.2, 11.5, 0.4,
             "Gap: Domain-specific tools are mature but siloed. Co-simulation frameworks exist but lack a reusable, modular engine architecture.",
             font_size=13, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# ── 5. RESEARCH APPROACH ──────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Research Approach", "Design Science Research (DSR)")

phases = [
    ("1", "Problem\nIdentification", "Gaps in cross-domain\nsimulation tools"),
    ("2", "Architecture\nDesign", "Three-layer modular\narchitecture with\nseparation of concerns"),
    ("3", "Framework\nDevelopment", "Python prototype,\ndomain modules,\nHELICS co-simulation"),
    ("4", "Empirical\nValidation", "Benchmarks, metrics,\ncoupled vs. uncoupled\ncomparison"),
]
x_start = 1.0
box_w = 2.6
gap = 0.35
for i, (num, title, desc) in enumerate(phases):
    x = x_start + i * (box_w + gap)
    _add_rect(sl, x + 0.85, 1.7, 0.7, 0.7, ACCENT, num, font_size=24, font_color=WHITE)
    _add_rect(sl, x, 2.6, box_w, 0.9, DARK_BG, title, font_size=15, font_color=WHITE)
    _add_textbox(sl, x, 3.65, box_w, 1.5, desc,
                 font_size=13, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    if i < 3:
        arrow_x = x + box_w + 0.02
        _add_textbox(sl, arrow_x, 2.7, 0.3, 0.7, "→",
                     font_size=28, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)

_add_textbox(sl, 1.0, 5.4, 11.0, 0.6,
             "DSR is suited to applied engineering where the goal is to build and evaluate a purposeful artifact.",
             font_size=14, color=MED_GRAY, alignment=PP_ALIGN.LEFT)

# ── 6. SYSTEM ARCHITECTURE OVERVIEW ───────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "System Architecture", "Three-layer modular design - Strict separation between generic infrastructure and domain-specific logic")

if ARCH_IMG.exists():
    image_width = Inches(7.0)
    image_height = Inches(4.4)
    sl.shapes.add_picture(str(ARCH_IMG), Inches(3.2), Inches(1.7), width=image_width)
    image_height_in = 4.4
else:
    image_height_in = 0.0

# ── 7. CORE ENGINE (LAYER 1) ──────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Layer 1: Core Simulation Engine", "Domain-independent infrastructure")

_add_bullet_list(sl, 0.7, 1.7, 5.8, 5.0, [
    "Agents: Autonomous entities with step() function invoked each tick",
    "Time Control: Discrete-event stepping with configurable Δt (100 ms → 15 min)",
    "Contracts: Modules implement initialize(), step(), generate_report()",
    "Collector: Generic metrics accumulation decoupled from domain reporting",
    "No assumptions about simulated entities - fully domain-agnostic",
], font_size=15, bold_prefix=True, bullet_color=LAYER1_CLR, line_spacing=1.5)

_add_rect(sl, 7.5, 1.8, 5.0, 0.6, DARK_BG, "Engine Contract Interface", font_size=16)
methods = ["initialize(config)", "step(time, dt)", "generate_report()", "get_state() → dict"]
for j, m in enumerate(methods):
    _add_rect(sl, 7.5, 2.55 + j * 0.65, 5.0, 0.55,
              LIGHT_GRAY, m, font_size=14, font_color=DARK_BG, bold=False)

_add_textbox(sl, 7.5, 5.2, 5.0, 0.5,
             "→ New domains plug in without modifying engine internals",
             font_size=13, color=ACCENT, bold=True)

# ── 8. DOMAIN MODULES (LAYER 2) ───────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Layer 2: Domain Modules", "Autonomous, independently runnable modules")

modules = [
    ("Energy", LAYER1_CLR,
     ["Solar PV, wind turbines",
      "Battery storage (BESS)",
      "Responsive loads + DR",
      "IEEE bus networks"]),
    ("Mobility", LAYER2_CLR,
     ["Agent-based vehicles",
      "A*-based routing",
      "Signalized intersections",
      "Adaptive signal control"]),
    ("Telecom", LAYER3_CLR,
     ["5G NR / gNBs",
      "Network slicing (eMBB,\nURLLC, mMTC)",
      "RB scheduling",
      "User mobility & handover"]),
]
for i, (name, clr, bullets) in enumerate(modules):
    x = 0.7 + i * 4.2
    _add_rect(sl, x, 1.7, 3.8, 0.7, clr, name + " Module", font_size=18)
    _add_bullet_list(sl, x + 0.2, 2.6, 3.4, 3.0, bullets,
                     font_size=14, bullet_color=clr, bold_prefix=False, line_spacing=1.6)

_add_textbox(sl, 0.7, 5.6, 11.5, 0.5,
             "Each module runs standalone or as a HELICS federate - only the use_helics flag changes; domain logic is identical.",
             font_size=14, color=MED_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

# ── 9. CO-SIMULATION LAYER (LAYER 3) ──────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Layer 3: Co-Simulation with HELICS",
             "Hierarchical Engine for Large-scale Infrastructure Co-Simulation")

_add_textbox(sl, 0.7, 1.6, 5.5, 0.5,
             "Why HELICS?", font_size=20, bold=True, color=DARK_BG)
_add_bullet_list(sl, 0.7, 2.2, 5.5, 2.5, [
    "Multi-domain by design (vs. Mosaik's smart-grid focus)",
    "Supports time-stepped + event-driven federates",
    "Validated at scale (100s of federates)",
    "Chosen over FMI/FMU (limited agent/network support)",
], font_size=14, bullet_color=ACCENT, bold_prefix=False, line_spacing=1.5)

_add_textbox(sl, 7.0, 1.6, 5.5, 0.5,
             "HELICS Manages:", font_size=20, bold=True, color=DARK_BG)
_add_bullet_list(sl, 7.0, 2.2, 5.5, 3.5, [
    "Federate Registration: Publications & subscriptions per module",
    "Time Synchronization: Grant-request protocol — a federate asks to advance, HELICS checks for earlier-time data, and only then grants the step",
    "Data Exchange: State variables flow through value federates (no API coupling)",
    "Broker: Single process coordinates up to 3+ federates",
], font_size=14, bullet_color=ACCENT, bold_prefix=True, line_spacing=1.5)

_add_rect(sl, 0.7, 5.0, 11.9, 0.7, LIGHT_GRAY,
          "Federation Pattern: Modules communicate exclusively through HELICS topics (publish-subscribe) - adding a new domain requires only new topics, not modifying existing modules",
          font_size=13, font_color=DARK_BG, bold=False)

# ── 10. SCENARIO OVERVIEW TABLE ───────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Simulation Scenarios Overview")

headers = ["Scenario", "Domain", "Description", "Duration", "Time Step", "Key Metrics"]
rows = [
    ["E1", "Energy", "Smart grid with renewables", "24 h", "15 min", "Voltage, power flow, curtailment"],
    ["E2", "Energy", "EV charging infrastructure", "36 h", "5 min", "Load profile, charging cost, V2G"],
    ["M1", "Mobility", "Urban traffic management", "3 h", "1 sec", "Travel time, delay, emissions"],
    ["T1", "Telecom", "5G slice resource allocation", "1 h", "100 ms", "Slice utilization, QoS, RB waste"],
    ["E2+M1", "Cross-domain", "Energy-Mobility integration", "3 h", "1 s / 5 min", "EV drive time, emissions, SOC"],
    ["M1+T1", "Cross-domain", "Mobility-Telecom integration", "3 h", "1 s / 100 ms", "Handovers, QoS, delay feedback"],
]
_add_table_slide(sl, headers, rows, left=0.5, top=1.6, width=12.3,
                 col_widths=[1.1, 1.4, 2.7, 1.1, 1.3, 4.7], font_size=13)

_add_textbox(sl, 0.5, 5.6, 12.0, 0.5,
             "All scenarios support both standalone execution and HELICS co-simulation mode.",
             font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ── 11. EVALUATION METHODOLOGY ────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Evaluation Methodology")

_add_textbox(sl, 0.7, 1.6, 5.5, 0.5,
             "Domain Model Validation", font_size=20, bold=True, color=DARK_BG)
_add_bullet_list(sl, 0.7, 2.2, 5.5, 3.5, [
    "E1: Voltage within IEEE 33-bus bounds; curtailment only when BESS saturated",
    "E2: Smart < Uncoord. peak; V2G reduces cost; ≥95% SOC target",
    "M1: Adaptive reduces travel time, delay, emissions vs. fixed",
    "T1: Dynamic improves utilization, reduces RB waste vs. static",
    "E2+M1: Coupled ≠ uncoupled (drive time, emissions, SOC deficit)",
    "M1+T1: Coupled ≠ uncoupled (handovers, QoS, delay feedback)",
], font_size=13, bullet_color=ACCENT, bold_prefix=False, line_spacing=1.5)

_add_textbox(sl, 7.0, 1.6, 5.5, 0.5,
             "Architectural Validation", font_size=20, bold=True, color=DARK_BG)
_add_bullet_list(sl, 7.0, 2.2, 5.5, 3.5, [
    "Domain Independence: Identical behavior in standalone & co-sim modes",
    "Interface Compliance: All agents implement engine contract",
    "Strategy Substitutability: Variants produce distinct outcomes under same input",
    "Time Synchronization: HELICS grants verified; data arrives before use",
], font_size=13, bullet_color=LAYER2_CLR, bold_prefix=True, line_spacing=1.5)

# ── 12. USER INTERFACE ────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "User Interface", "")
poster = str(UI_IMG) if UI_IMG.exists() else None
sl.shapes.add_movie(
    str(DEMO_VIDEO),
    Inches(1.67), Inches(1.2),
    width=Inches(10.0), height=Inches(5.625),
    poster_frame_image=poster,
    mime_type="video/mp4",
)

# ── 13. TECHNICAL CHALLENGES ──────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Technical Challenges")

# Add M1 scaling performance plot
M1_PLOT = ROOT / "figs" / "m1_scaling_performance.png"
sl.shapes.add_picture(str(M1_PLOT), Inches(1.1), Inches(1.2), width=Inches(11))

# Challenge cards below plot
challenges = [
    (LAYER2_CLR, "M1 Scalability (27 M agent-steps)",
     "A* computed once per vehicle at departure & cached;\n"
     "O(N) queue scan acceptable at N=2,500;\n"
     "E2+M1 wall-clock ≈ 25 s on dev machine"),
    (LAYER3_CLR, "Cross-Domain Coordinate Mismatch",
     "E2: IEEE 13-node indices  ↔  M1: (x,y) grid tuples;\n"
     "_station_by_pos dict → O(1) lookup;\n"
     "M1+T1: T1 area_size rescaled to 5,000 m"),
]
for i, (clr, title, desc) in enumerate(challenges):
    x = 0.7 + i * 6.1
    y = 4.9
    _add_rect(sl, x, y, 5.7, 0.55, clr, title, font_size=13)
    _add_textbox(sl, x + 0.1, y + 0.65, 5.5, 1.3, desc,
                 font_size=12, color=BLACK, alignment=PP_ALIGN.LEFT)

# ── 14. RESULTS SUMMARY ───────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Results Summary", "Key findings across all six scenarios")

# Left side: Summary cards (3 cards in a single vertical column)
summary_cards = [
    (LAYER3_CLR, "T1 - 5G Slicing",
     "Dynamic slicing improves broadband QoS while revealing a clear fairness trade-off for IoT traffic."),
    (ACCENT, "E2+M1 - Cross-domain",
     "Coupled EV routing adds real congestion effects that are invisible when energy and mobility are simulated separately."),
    (ACCENT2, "M1+T1 - Cross-domain",
     "Mobility and telecom feedback jointly degrade signal control and increase intersection delay under congestion."),
]
for i, (clr, ttl, desc) in enumerate(summary_cards):
    x = 0.5
    y = 1.9 + i * 1.75
    _add_rect(sl, x, y, 5.0, 0.42, clr, ttl, font_size=14)
    _add_textbox(sl, x + 0.12, y + 0.52, 5.25, 1.05, desc,
                 font_size=12, color=BLACK, alignment=PP_ALIGN.LEFT)

# Right side: Energy and Mobility charts stacked vertically
energy_chart = ROOT / "figs" / "energy_domain_summary.png"
mobility_chart = ROOT / "figs" / "mobility_domain_summary.png"

if energy_chart.exists():
    sl.shapes.add_picture(str(energy_chart), Inches(6.2), Inches(1.5), width=Inches(6.5))
else:
    _add_rect(sl, 6.2, 1.5, 6.5, 2.5, LIGHT_GRAY,
              "[Energy Domain Summary]",
              font_size=12, font_color=MED_GRAY, bold=False)

if mobility_chart.exists():
    sl.shapes.add_picture(str(mobility_chart), Inches(6.2), Inches(4.0), width=Inches(6.5))
else:
    _add_rect(sl, 6.2, 4.0, 6.5, 2.5, LIGHT_GRAY,
              "[Mobility Domain Summary]",
              font_size=12, font_color=MED_GRAY, bold=False)

# ── 15. MAIN CONTRIBUTIONS ────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Main Contributions")

_add_bullet_list(sl, 0.7, 1.7, 11.9, 5.5, [
    "Domain-Agnostic Simulation Engine: BaseFederate ABC with 3-method contract; transparent standalone/co-sim dual mode via use_helics flag; no domain knowledge in engine internals",
    "Multi-Domain Scenario Framework: 6 autonomous modules spanning energy (E1, E2), mobility (M1), telecommunications (T1), and two cross-domain integrations (E2+M1, M1+T1)",
    "Quantification of Cross-Domain Interactions: EV traffic +21.5% CO₂ & -7.0% SoC targets (E2+M1); vehicular RF clustering raises intersection delay +80.9% via T1→M1 feedback (M1+T1)",
    "Strategy Substitutability via Design Patterns: ChargingStrategy, adaptive_signals, SlicingStrategy enums enable comparative evaluation without code duplication - validates engine contract",
    "Reproducible & Accessible Prototype: Streamlit dashboard for no-code execution; deterministic results from any seed; public repository with requirements.txt for full reproducibility",
], font_size=16, bold_prefix=True, bullet_color=ACCENT, line_spacing=1.45)

# ── 16. OBJECTIVES ACHIEVEMENT ────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Objectives Achievement")

headers_obj = ["Objective", "Description", "Status"]
rows_obj = [
    ["Obj. 1", "Domain-agnostic simulation architecture (BaseFederate)", "ACHIEVED"],
    ["Obj. 2", "Domain-specific modules (energy, mobility, telecom)", "ACHIEVED"],
    ["Obj. 3", "Multi-domain co-simulation via HELICS", "ACHIEVED"],
    ["Obj. 4", "Reproducibility and configurability", "ACHIEVED"],
    ["Obj. 5", "Validation through practical scenarios (≥2 domains)", "ACHIEVED"],
    ["Obj. 6", "Urban Management Platform integration (Ubiwhere UBP)", "PARTIAL"],
]
n_rows = len(rows_obj) + 1
n_cols = 3
table_shape = sl.shapes.add_table(n_rows, n_cols,
                                   Inches(0.7), Inches(1.6),
                                   Inches(11.9), Inches(0.5 * n_rows))
table = table_shape.table
for cw, w in zip(range(3), [1.4, 7.5, 1.8]):
    table.columns[cw].width = Inches(w)

for j, h in enumerate(headers_obj):
    cell = table.cell(0, j)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = DARK_BG
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14); p.font.bold = True
        p.font.color.rgb = WHITE; p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

for i, row in enumerate(rows_obj):
    for j, val in enumerate(row):
        cell = table.cell(i + 1, j)
        cell.text = val
        if i % 2 == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GRAY
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13); p.font.name = "Calibri"
            p.font.color.rgb = BLACK; p.alignment = PP_ALIGN.CENTER
            if j == 2:
                if val == "ACHIEVED":
                    p.font.color.rgb = GREEN; p.font.bold = True
                    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xD4, 0xED, 0xDA)
                elif val == "PARTIAL":
                    p.font.color.rgb = ORANGE; p.font.bold = True
                    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xCD)
            if j == 1:
                p.alignment = PP_ALIGN.LEFT

_add_rect(sl, 0.7, 5.45, 11.9, 0.75, LIGHT_GRAY,
          "Ubiwhere Alignment: Three-layer architecture with pub-sub patterns directly enables integration into Ubiwhere UBP. "
          "Modular design ensures simulation engines can feed real-time what-if scenarios to urban management dashboards. ",
          font_size=12, font_color=DARK_BG, bold=False)

# ── 17. LIMITATIONS ───────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Limitations", "Self-imposed design choices to maintain prototype simplicity")

_add_textbox(sl, 0.7, 1.95, 5.7, 0.4,
             "Model Fidelity", font_size=18, bold=True, color=DARK_BG)
_add_bullet_list(sl, 0.7, 2.45, 5.7, 3.5, [
    "E1 voltage model: linear approximation, not AC power-flow",
    "E2 perfect information: controller assumes known departure time & battery state",
], font_size=14, bullet_color=ORANGE, bold_prefix=False, line_spacing=1.5)

_add_textbox(sl, 7.0, 1.95, 5.7, 0.4,
             "Architecture & Scale", font_size=18, bold=True, color=DARK_BG)
_add_bullet_list(sl, 7.0, 2.45, 5.7, 3.5, [
    "Static routing: M1 routes computed once at departure, not updated for changing congestion",
    "Domain coverage: only 3 verticals (water, waste, safety not addressed)",
], font_size=14, bullet_color=RED_LIGHT, bold_prefix=False, line_spacing=1.5)

_add_rect(sl, 0.7, 6.05, 11.9, 0.65, LIGHT_GRAY,
          "All limitations are deliberate trade-offs: chosen to keep the prototype manageable while preserving architectural soundness. "
          "The three-layer design is fully extensible - each limitation lifts independently without restructuring.",
          font_size=12, font_color=DARK_BG, bold=False)

# ── 18. FUTURE WORK ───────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Future Work")

_add_textbox(sl, 0.7, 1.6, 5.7, 0.4,
             "Near-Term Extensions", font_size=18, bold=True, color=DARK_BG)
_add_bullet_list(sl, 0.7, 2.1, 5.7, 4.0, [
    "UBP integration: connectors, data-model adapter (Obj. 6 completion)",
    "Full cross-domain scenario: M2+E2+T1 (autonomous vehicles, EV charging, 5G slices)",
    "AC power-flow solver for E1 (pandapower / OpenDSS) - eliminates voltage over-estimation",
    "Dynamic vehicle rerouting in M1 (real-time congestion feedback for connected vehicles)",
    "M2+T1 & E1+T2 cross-domain scenarios",
], font_size=16, bullet_color=ACCENT, bold_prefix=False, line_spacing=1.45)

_add_textbox(sl, 7.0, 1.6, 5.7, 0.4,
             "Long-Term Research Directions", font_size=18, bold=True, color=DARK_BG)
_add_bullet_list(sl, 7.0, 2.1, 5.7, 4.0, [
    "More indicators: Improve the quality of validation metrics",
    "Scenario M2: Autonomous vehicle integration (mixed autonomy, platoon, V2X)",
    "Scenario T2: Dense IoT network (thousands of low-power devices, collision model)",
    "City-scale deployment: spatial indexing, parallel agent stepping, distributed HELICS federation",
    "Additional urban verticals: water distribution, waste management, public safety",
], font_size=16, bullet_color=ACCENT2, bold_prefix=False, line_spacing=1.45)

# ── 19. FINAL REMARKS & QUESTIONS ─────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, DARK_BG)
_slide_title(sl, "Final Remarks", dark=True)

_add_bullet_list(sl, 0.8, 1.75, 11.5, 4.0, [
    "Three-layer architecture achieves strict separation: engine ↔ domain modules ↔ HELICS orchestration",
    "Dual-mode execution (standalone / co-sim) validated - domain logic identical under both modes",
    "Cross-domain coupling reveals effects invisible to siloed simulation:",
    "   → E2+M1: 500 EVs increase fleet CO₂ by 21.5%",
    "   → M1+T1: vehicular RF clustering raises intersection delay by 80.9% via T1→M1 feedback",
    "Extensibility evaluated: E2+M1 reuses E2 & M1 components with zero modifications to source files",
    "Streamlit dashboard enables no-code scenario execution and strategy comparison",
], font_size=17, color=WHITE, bullet_color=ACCENT, bold_prefix=False, line_spacing=1.5)

# accent line
line2 = sl.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(5.85), Inches(6.3), Inches(0.06))
line2.fill.solid(); line2.fill.fore_color.rgb = ACCENT; line2.line.fill.background()

_add_textbox(sl, 1.0, 6.05, 11.3, 0.8,
             "Thank you - Questions?",
             font_size=30, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# ── 20. BASEFEDERATE CORE ENGINE ──────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "BaseFederate: Core Engine", "engine/base.py - domain-agnostic simulation infrastructure")

_add_bullet_list(sl, 0.7, 1.7, 5.8, 4.8, [
    "Constructor: name, use_helics, time_step, sim_duration",
    "Abstract method 1: initialize_components() - create agents",
    "Abstract method 2: run_simulation() - main stepping loop",
    "Abstract method 3: generate_report() → Dict - metrics",
    "Concrete: setup_federate() - HELICS init or graceful fallback",
    "Concrete: advance_time(t) - HELICS grant or arithmetic step",
    "Dual-mode: use_helics flag only; all domain logic identical",
    "HELICS failure → catches ImportError/runtime, logs warning, continues standalone",
], font_size=14, bold_prefix=True, bullet_color=LAYER1_CLR, line_spacing=1.4)

_add_rect(sl, 7.5, 1.8, 5.3, 0.55, DARK_BG, "BaseFederate (ABC)", font_size=15)
contract = [
    "initialize_components(self)",
    "run_simulation(self)",
    "generate_report(self) → Dict",
    "─────────────────────────",
    "setup_federate()",
    "advance_time(current_t)",
]
for j, m in enumerate(contract):
    bg = LAYER1_CLR if j < 3 else LIGHT_GRAY
    fc = WHITE if j < 3 else DARK_BG
    _add_rect(sl, 7.5, 2.5 + j * 0.6, 5.3, 0.52,
              bg, m, font_size=12, font_color=fc, bold=(j < 3))

# ── 21. RESULTS E1 ────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Results: E1 - Smart Grid with Renewables")

headers_e1 = ["Metric", "Value", "Unit"]
rows_e1 = [
    ["Total renewable generation", "13,744", "kWh"],
    ["  of which solar PV", "2,162", "kWh"],
    ["  of which wind", "11,582", "kWh"],
    ["Total load consumed", "22,646", "kWh"],
    ["Renewable penetration", "60.7", "%"],
    ["Total curtailment", "1,365", "kWh"],
    ["Curtailment rate", "9.9", "% of generation"],
    ["Peak load", "1,262", "kW"],
    ["Estimated DR load reduction", "147.7", "kW (10.5%)"],
    ["Average bus voltage", "0.998", "p.u."],
    ["Minimum bus voltage", "0.969", "p.u."],
    ["Maximum bus voltage", "1.303", "p.u.  ⚠"],
    ["Voltage violations (steps)", "49 / 96", "-"],
    ["Curtailment validation", "PASS", "-"],
]
_add_table_slide(sl, headers_e1, rows_e1, left=0.7, top=1.6, width=7.5,
                 col_widths=[4.2, 1.8, 1.5], font_size=13)

_add_rect(sl, 8.5, 1.6, 4.5, 1.1, LIGHT_GRAY,
          "Wind = 84.3% of renewables;\n60.7% penetration exceeds base load;\ncurtailment in 18 / 96 steps (BESS full)",
          font_size=13, font_color=DARK_BG, bold=False)
_add_rect(sl, 8.5, 2.9, 4.5, 1.4, RGBColor(0xFF, 0xF3, 0xCD),
          "⚠  Max voltage 1.303 p.u. exceeds\nANSI C84.1 limit (1.05 p.u.).\nLinear voltage-drop model lacks\nreactive-power regulation.",
          font_size=13, font_color=RGBColor(0x85, 0x60, 0x04), bold=False)
_add_rect(sl, 8.5, 4.5, 4.5, 0.9, DARK_BG,
          "Future work: replace linear model\nwith Newton-Raphson AC power-flow\n(pandapower / OpenDSS)",
          font_size=12, font_color=WHITE, bold=False)

# ── 22. RESULTS E2 ────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Results: E2 - EV Charging Infrastructure")

headers_e2 = ["Metric", "Uncoordinated", "Smart", "V2G"]
rows_e2 = [
    ["Peak load (kW)", "2,512", "2,450", "2,450"],
    ["Total energy charged (kWh)", "3,112", "3,731", "4,322"],
    ["Total charging cost (USD)", "$456.44", "$511.96", "$327.88"],
    ["Avg. cost per vehicle (USD)", "$4.56", "$5.12", "$3.28"],
    ["Avg. final SoC", "0.81", "0.90", "0.88"],
    ["Max transformer loading (%)", "117.3", "114.4", "114.4"],
    ["Transformer overloads (events)", "131", "135", "88"],
    ["V2G energy provided (kWh)", "0.0", "0.0", "66.1"],
    ["V2G revenue (USD)", "$0.00", "$0.00", "$171.20"],
    ["Vehicles meeting SoC target", "100", "100", "97"],
    ["Load factor", "0.60", "0.62", "0.63"],
]
_add_table_slide(sl, headers_e2, rows_e2, left=0.5, top=1.6, width=12.3,
                 col_widths=[4.0, 2.2, 2.0, 2.1], font_size=13)

key_insight = (
    "V2G lowest cost ($327.88) via $171.20 export revenue  ·  "
    "Smart reduces peak but costs more (urgency-first ignores price)  ·  "
    "3 V2G vehicles miss SoC target (discharged before re-charge)"
)
_add_rect(sl, 0.5, 5.95, 12.3, 0.7, DARK_BG, key_insight,
          font_size=13, font_color=WHITE, bold=False)

# ── 23. RESULTS M1 ────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Results: M1 - Urban Traffic Management")

headers_m1 = ["Metric", "Fixed-Time", "Adaptive", "Change"]
rows_m1 = [
    ["Completed vehicles", "2,500", "2,500", "-"],
    ["Completion rate (%)", "100.0", "100.0", "-"],
    ["Avg. travel time (s)", "583.9", "540.7", "-7.4%"],
    ["Avg. delay (s)", "66.0", "22.8", "-65.5%"],
    ["Total system delay (s)", "164,987", "56,942", "-65.5%"],
    ["Total CO₂ emissions (kg)", "3,081", "2,831", "-8.1%"],
    ["Emissions per vehicle (g)", "1,232", "1,132", "-8.1%"],
    ["Max. queue length", "17", "13", "-23.5%"],
    ["Avg. queue length", "0.66", "0.26", "-60.6%"],
    ["Throughput (veh/h)", "833", "833", "-"],
]
_add_table_slide(sl, headers_m1, rows_m1, left=0.6, top=1.6, width=11.5,
                 col_widths=[3.8, 2.1, 2.1, 1.8], font_size=13)

_add_rect(sl, 0.6, 5.95, 11.5, 0.7, LAYER2_CLR,
          "Adaptive signals cut delay by 65.5% (66.0 s → 22.8 s) and CO₂ by 8.1%  ·  "
          "Throughput unchanged - benefit is in efficiency, not volume  ·  "
          "Queue-responsive green formula prevents build-up at saturated approaches",
          font_size=13, font_color=WHITE, bold=False)

# ── 24. RESULTS T1 ────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Results: T1 - 5G Slice Resource Allocation")

headers_t1 = ["Metric", "Static", "Dynamic", "Change"]
rows_t1 = [
    ["QoS satisfaction - eMBB (%)", "8.6", "49.8", "+479%"],
    ["QoS satisfaction - URLLC (%)", "55.7", "44.8", "-19.6%"],
    ["QoS satisfaction - mMTC (%)", "87.1", "48.4", "-44.4%"],
    ["RB utilisation - eMBB (%)", "71.8", "79.2", "+10.3%"],
    ["RB utilisation - URLLC (%)", "69.2", "79.6", "+15.0%"],
    ["RB utilisation - mMTC (%)", "53.9", "53.3", "-1.1%"],
    ["Overall RB utilisation (%)", "67.5", "75.3", "+11.6%"],
    ["Resource waste (%)", "32.5", "24.7", "-24.0%"],
    ["Handover success rate (%)", "95.0", "95.0", "-"],
]
_add_table_slide(sl, headers_t1, rows_t1, left=0.6, top=1.6, width=11.5,
                 col_widths=[3.8, 1.8, 1.8, 1.8], font_size=13)

_add_rect(sl, 0.6, 5.6, 5.5, 1.1, LIGHT_GRAY,
          "Static: 50% RBs to eMBB → only 50 RBs for 100 users * 2 RB req.\n"
          "→ eMBB QoS collapses to 8.6%  (demand 200 RBs vs. 50 available)",
          font_size=13, font_color=DARK_BG, bold=False)
_add_rect(sl, 6.3, 5.6, 5.8, 1.1, ORANGE,
          "Trade-off: Dynamic raises eMBB +479% but mMTC drops -44.4%.\n"
          "Neither strategy achieves uniform QoS - constraint-based\noptimisation needed for heterogeneous SLA targets.",
          font_size=13, font_color=WHITE, bold=False)

# ── 25. RESULTS E2+M1 ─────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Results: E2+M1 - Energy-Mobility Cross-Domain")

headers_em = ["Metric", "Uncoupled", "Coupled", "Change"]
rows_em = [
    ["- Mobility metrics -", "", "", ""],
    ["Avg. travel time (s)", "536.2", "536.4", "+0.04%"],
    ["Avg. delay (s)", "22.5", "22.6", "+0.4%"],
    ["Total emissions (kg CO₂)", "2,246", "2,729", "+21.5%"],
    ["- EV travel metrics -", "", "", ""],
    ["Avg. EV drive time (s)", "0.5", "366.5", "-"],
    ["Avg. EV total-to-station (s)", "3,069", "3,154", "+2.8%"],
    ["EVs that reached station", "155", "150", "-3.2%"],
    ["- Charging metrics -", "", "", ""],
    ["EVs meeting SoC target", "71", "66", "-7.0%"],
    ["Total energy charged (kWh)", "3,189", "3,127", "-1.9%"],
    ["Total charging cost (USD)", "$255.15", "$250.17", "-2.0%"],
    ["Peak grid load (kW)", "2,107", "2,107", "-"],
]
_add_table_slide(sl, headers_em, rows_em, left=0.5, top=1.6, width=12.3,
                 col_widths=[4.0, 2.0, 2.0, 1.8], font_size=12)

_add_rect(sl, 0.5, 6.15, 12.3, 0.65, ACCENT,
          "Cross-domain effect: 500 EVs add +21.5% CO₂ and delay 5/500 EVs from reaching SoC target - "
          "invisible to domain-specific simulation of either energy or mobility alone",
          font_size=13, font_color=WHITE, bold=True)

# ── 26. RESULTS M1+T1 ─────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Results: M1+T1 - Mobility-Telecom Cross-Domain")

headers_mt = ["Metric", "Uncoupled", "Coupled", "Change"]
rows_mt = [
    ["- Telecommunications -", "", "", ""],
    ["QoS satisfaction - eMBB (%)", "39.0", "15.4", "-60.5%"],
    ["QoS satisfaction - URLLC (%)", "39.4", "19.2", "-51.3%"],
    ["QoS satisfaction - mMTC (%)", "40.5", "16.9", "-58.3%"],
    ["Overall RB utilisation (%)", "73.8", "60.1", "-18.6%"],
    ["Total handover events", "134,843", "156,168", "+15.8%"],
    ["Per-gNB load std dev (users)", "15.6", "7.0", "-55.1%"],
    ["- Mobility (all vehicles) -", "", "", ""],
    ["Avg. delay (s)", "22.7", "41.0", "+80.9%"],
    ["Signal degradation events", "0", "25,915", "-"],
    ["Intersections affected", "0", "14 / 25", "-"],
    ["Total CO₂ emissions (kg)", "2,827", "2,933", "+3.7%"],
]
_add_table_slide(sl, headers_mt, rows_mt, left=0.5, top=1.6, width=12.3,
                 col_widths=[3.8, 2.1, 2.1, 1.8], font_size=12)

_add_rect(sl, 0.5, 6.1, 12.3, 0.7, DARK_BG,
          "T1→M1 causal chain: vehicular RF clustering → URLLC QoS <80% → 14/25 intersections revert to fixed-time → "
          "+80.9% intersection delay  ·  Only observable through coupled simulation",
          font_size=13, font_color=WHITE, bold=True)

# ── 27. ENERGY SCENARIOS (E1 + E2) ────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Energy Scenarios: E1 & E2")

_add_rect(sl, 0.7, 1.6, 5.7, 0.55, LAYER1_CLR,
          "E1: Smart Grid with Renewable Integration", font_size=15)
_add_bullet_list(sl, 0.7, 2.3, 5.7, 3.0, [
    "IEEE 33-bus radial distribution network",
    "50 solar PV (5-10 kW) + 3 wind turbines (500 kW)",
    "5 BESS units (50 kWh / 25 kW, η = 0.92)",
    "800 residential loads with demand response",
    "Voltage model: simplified linear drop (ANSI C84.1)",
    "Curtailment when batteries full + surplus remains",
], font_size=13, bullet_color=LAYER1_CLR, bold_prefix=False, line_spacing=1.4)

_add_rect(sl, 6.8, 1.6, 5.8, 0.55, LAYER1_CLR,
          "E2: EV Charging Infrastructure", font_size=15)
_add_bullet_list(sl, 6.8, 2.3, 5.8, 3.0, [
    "100 EVs on IEEE 13-node feeder (2500 kW capacity)",
    "80 residential (L2, 7.2 kW) + 10 DC-fast (50 kW) ports",
    "Three strategies: Uncoordinated → Smart → V2G",
    "Smart: priority by urgency U = E_need / T_available",
    "V2G: bidirectional flow during peak (17-21h), 20% premium",
    "Time-of-Use tariff: $0.08 / $0.12 / $0.20 per kWh",
], font_size=13, bullet_color=LAYER1_CLR, bold_prefix=False, line_spacing=1.4)

_add_rect(sl, 0.7, 5.6, 11.9, 0.6, LIGHT_GRAY,
          "Strategy Pattern: ChargingStrategy enum selects UNCOORDINATED / SMART / V2G at runtime - comparative evaluation without code duplication",
          font_size=13, font_color=DARK_BG, bold=False)

# ── 28. MOBILITY SCENARIO M1 ──────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Mobility Scenario: M1", "Urban Traffic Congestion Management")

_add_bullet_list(sl, 0.7, 1.7, 5.8, 4.5, [
    "5*5 km grid, 25 signalized intersections (1250 m spacing)",
    "2,500 agent-based vehicles, 50 km/h link speed",
    "Routing: A* algorithm with Manhattan heuristic + stochastic diversity",
    "1-second resolution, 3 hours (10,800 steps)",
], font_size=15, bullet_color=LAYER2_CLR, bold_prefix=False, line_spacing=1.6)

_add_textbox(sl, 7.0, 1.6, 5.5, 0.5,
             "Signal Control Comparison", font_size=18, bold=True, color=DARK_BG)
_add_rect(sl, 7.0, 2.2, 5.5, 0.55, MED_GRAY,
          "Fixed-Time: Equal green for both phases", font_size=14)
_add_rect(sl, 7.0, 2.9, 5.5, 0.55, LAYER2_CLR,
          "Adaptive: Green ∝ measured queue demand", font_size=14)
_add_textbox(sl, 7.0, 3.7, 5.5, 0.8,
             "T_green = T_min + (T_max - T_min) · Q_phase / Q_total\nT_min = 15s, T_max = 90s",
             font_size=14, color=DARK_BG, alignment=PP_ALIGN.LEFT)
_add_textbox(sl, 7.0, 4.6, 5.5, 0.4,
             "Emissions Model:", font_size=16, bold=True, color=DARK_BG)
_add_bullet_list(sl, 7.0, 5.0, 5.5, 1.5, [
    "Idle (v < 1 m/s): 2.31 g/s",
    "Moving: 0.15 g/m",
], font_size=14, bullet_color=LAYER2_CLR, bold_prefix=False, line_spacing=1.4)

# ── 29. TELECOM SCENARIO T1 ───────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Telecom Scenario: T1", "5G Slice Resource Allocation")

_add_bullet_list(sl, 0.7, 1.7, 5.8, 4.5, [
    "3 gNBs, 100 RBs each (20 MHz, 15 kHz SCS)",
    "200 users: 100 eMBB + 40 URLLC + 60 mMTC",
    "3GPP UMa path loss model (TR 38.901)",
    "Random Waypoint mobility (0.5-2.0 m/s)",
    "Handover: 3 dB hysteresis, 1s cooldown, 95% success",
    "mMTC: 30% base activity + periodic bursts (50s cycle)",
], font_size=15, bullet_color=LAYER3_CLR, bold_prefix=False, line_spacing=1.5)

_add_textbox(sl, 7.0, 1.6, 5.5, 0.5,
             "Slicing Strategies", font_size=18, bold=True, color=DARK_BG)
_add_rect(sl, 7.0, 2.2, 5.5, 0.55, MED_GRAY,
          "Static: Fixed 50% eMBB / 30% URLLC / 20% mMTC", font_size=13)
_add_rect(sl, 7.0, 2.9, 5.5, 0.55, LAYER3_CLR,
          "Dynamic: RBs ∝ instantaneous demand (floor = 5 RBs)", font_size=13)
_add_textbox(sl, 7.0, 3.7, 5.5, 1.0,
             "QoS satisfied when allocated RBs ≥ requirement:\n• eMBB / URLLC: 2 RBs\n• mMTC: 1 RB",
             font_size=14, color=DARK_BG)
_add_textbox(sl, 7.0, 4.9, 5.5, 0.8,
             "Analogous to M1's adaptive signals:\ncapacity allocated proportional to demand\nrather than fixed a priori.",
             font_size=13, color=MED_GRAY)

# ── 30. CROSS-DOMAIN E2+M1 ────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Cross-Domain: E2+M1", "Energy-Mobility Integration via HELICS")

_add_bullet_list(sl, 0.7, 1.7, 5.8, 3.5, [
    "500 EVs + 2000 background vehicles on 5*5 grid",
    "13 charging stations (5 DC-fast + 8 Level 2, 52 ports)",
    "Multi-rate: M1 @ 1s base tick, E2 charging logic every 300 ticks (5 min)",
    "EV lifecycle: DRIVING_TO_STATION → CHARGING → DRIVING_BACK → DONE",
    "Smart charging allocation reuses E2's urgency-based strategy",
], font_size=14, bullet_color=ACCENT, bold_prefix=False, line_spacing=1.5)

_add_textbox(sl, 7.0, 1.6, 5.5, 0.5,
             "Coupled vs. Uncoupled", font_size=18, bold=True, color=DARK_BG)
_add_rect(sl, 7.0, 2.2, 5.5, 1.0, LAYER2_CLR,
          "Coupled: EVs navigate traffic,\nwait at signals, cause congestion", font_size=14)
_add_rect(sl, 7.0, 3.4, 5.5, 1.0, MED_GRAY,
          "Uncoupled: EVs teleport to station\n(treats domains as independent)", font_size=14)
_add_textbox(sl, 7.0, 4.7, 5.5, 0.4,
             "Coupling-Specific Metrics:", font_size=16, bold=True, color=DARK_BG)
_add_bullet_list(sl, 7.0, 5.1, 5.5, 2.0, [
    "Avg EV Drive Time (≈0 uncoupled vs. real delay coupled)",
    "Emission Uplift (additional CO₂ from EV driving)",
    "SOC Target Deficit (fewer EVs reach target due to travel delays)",
], font_size=13, bullet_color=ACCENT, bold_prefix=False, line_spacing=1.4)

# ── 31. CROSS-DOMAIN M1+T1 ────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(sl, WHITE)
_slide_title(sl, "Cross-Domain: M1+T1", "Mobility-Telecommunications Integration via HELICS")

_add_bullet_list(sl, 0.7, 1.7, 5.8, 3.5, [
    "2,300 background + 200 ConnectedVehicle agents on 5*5 grid",
    "200 connected vehicles are also T1 users; positions drive RF model",
    "ConnectedVehicle.sync_position() maps grid node → T1 metres",
    "T1 path-loss, handover & RB allocation run once per M1 tick (1 s)",
    "QoS feedback: gNBs below 80% URLLC → revert signals to fixed-time",
], font_size=14, bullet_color=ACCENT, bold_prefix=False, line_spacing=1.5)

_add_textbox(sl, 7.0, 1.6, 5.5, 0.5,
             "T1→M1 Feedback Loop", font_size=18, bold=True, color=DARK_BG)
_add_rect(sl, 7.0, 2.2, 5.5, 1.1, LAYER3_CLR,
          "Coupled: Vehicular RF clustering\ndegrades URLLC QoS → suppresses\nadaptive signal control", font_size=13)
_add_rect(sl, 7.0, 3.45, 5.5, 0.9, MED_GRAY,
          "Uncoupled: 200 T1 users use RWP;\nQoS feedback path disabled", font_size=13)
_add_textbox(sl, 7.0, 4.55, 5.5, 0.4,
             "Coupling-Specific Metrics:", font_size=16, bold=True, color=DARK_BG)
_add_bullet_list(sl, 7.0, 5.0, 5.5, 2.0, [
    "Handover rate uplift (vehicles cross cells faster than pedestrians)",
    "QoS degradation in coupled mode (burst loading at red-light releases)",
    "Signal degradation events & intersection delay uplift (T1→M1 effect)",
], font_size=13, bullet_color=ACCENT, bold_prefix=False, line_spacing=1.4)

# ======================================================================
#  FOOTERS
# ======================================================================
total = len(prs.slides)
for idx, slide in enumerate(prs.slides, start=1):
    _add_footer(slide, idx, total, dark=_is_dark_background(slide))

notes_map = _load_presenter_notes(NOTES_FILE)
_apply_presenter_notes(prs, notes_map)

# ── Save ───────────────────────────────────────────────────────────────
prs.save(str(OUT))
print(f"✅  Presentation saved to: {OUT}")
print(f"    Total slides: {total}")
